import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation


class DocumentProcessor:
    """Extracts clean text from PDF, DOCX, and PPTX files."""

    @staticmethod
    def extract_from_pdf(file_bytes: bytes) -> str:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                pages.append(text.strip())
        doc.close()
        return "\n\n".join(pages)

    @staticmethod
    def extract_from_docx(file_bytes: bytes) -> str:
        document = docx.Document(io.BytesIO(file_bytes))
        paragraphs = []
        for para in document.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # Also extract from tables
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        paragraphs.append(text)
        return "\n\n".join(paragraphs)

    @staticmethod
    def extract_from_pptx(file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))
        return "\n\n".join(slides)

    @classmethod
    def extract(cls, file_bytes: bytes, file_type: str) -> str:
        """Route to the correct extractor based on file type."""
        extractors = {
            "pdf": cls.extract_from_pdf,
            "docx": cls.extract_from_docx,
            "pptx": cls.extract_from_pptx,
        }
        extractor = extractors.get(file_type)
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        return extractor(file_bytes)
